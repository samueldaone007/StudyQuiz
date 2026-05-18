"""
Utility functions for text summarization, file parsing, and quiz generation.
"""
import re
import json
import random
import requests
import os
from typing import List, Dict, Tuple, Optional
from django.conf import settings

# File parsing imports
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def parse_docx_file(file_path: str) -> str:
    """Extract text from a DOCX file."""
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx is not installed. Run: pip install python-docx")
    
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        raise ValueError(f"Error parsing DOCX file: {str(e)}")


def parse_pdf_file(file_path: str) -> str:
    """Extract text from a PDF file."""
    if not PDF_AVAILABLE:
        raise ImportError("PyPDF2 is not installed. Run: pip install PyPDF2")
    
    try:
        text = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return '\n'.join(text)
    except Exception as e:
        raise ValueError(f"Error parsing PDF file: {str(e)}")


def parse_pptx_file(file_path: str) -> str:
    """
    Extract text from a PowerPoint file, preserving slide structure.

    Each slide is labelled with its number and title so the summariser
    can distinguish headings (Step 1, Step 2) from body content,
    examples, and code output blocks.
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is not installed. Run: pip install python-pptx")

    CODE_OUTPUT_MARKERS = re.compile(
        r'^\s*(>>>|output|in \[\d+\]|out\[\d+\]|#\s*(install|import)|'
        r'pip\d*\s+install|python\d*\s+-m|import\s+\w+|from\s+\w+\s+import)',
        re.IGNORECASE | re.MULTILINE,
    )

    try:
        prs = Presentation(file_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            title = ''
            body_shapes = []

            for shape in slide.shapes:
                if not hasattr(shape, 'text') or not shape.text.strip():
                    continue

                shape_text = shape.text.strip()

                if shape.shape_type == 13:
                    continue
                try:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                        ph_idx = shape.placeholder_format.idx
                        if ph_idx == 0 or ph_idx == 1:
                            title = shape_text
                            continue
                except Exception:
                    pass

                if CODE_OUTPUT_MARKERS.search(shape_text):
                    continue

                lines = shape_text.split('\n')
                prose_lines = [
                    l for l in lines
                    if re.search(r'[a-zA-Z]{4,}', l) and not re.match(
                        r'^\s*(>>>|#|pip|import|from\s+\w+|print\(|for |if |def |class )', l
                    )
                ]
                if not prose_lines:
                    continue

                body_shapes.append('\n'.join(prose_lines))

            if not title and not body_shapes:
                continue

            slide_block = f"[Slide {slide_num}: {title}]\n" + '\n'.join(body_shapes)
            slides_text.append(slide_block)

        return '\n\n'.join(slides_text)

    except Exception as e:
        raise ValueError(f"Error parsing PowerPoint file: {str(e)}")


def parse_uploaded_file(file_path: str) -> str:
    """
    Parse an uploaded file and extract text content.
    
    Args:
        file_path: Path to the uploaded file
    
    Returns:
        Extracted text content
    """
    if not os.path.exists(file_path):
        raise ValueError("File does not exist")
    
    # Get file extension
    _, ext = os.path.splitext(file_path.lower())
    
    if ext == '.txt':
        # Plain text file
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
    
    elif ext == '.docx':
        return parse_docx_file(file_path)
    
    elif ext == '.pdf':
        return parse_pdf_file(file_path)
    
    elif ext in ['.pptx', '.ppt']:
        return parse_pptx_file(file_path)
    
    else:
        raise ValueError(f"Unsupported file format: {ext}. Supported formats: .txt, .docx, .pdf, .pptx, .ppt")


def count_words(text: str) -> int:
    """Count words in text."""
    words = re.findall(r'\b\w+\b', text)
    return len(words)


def _is_technical_content(text: str) -> bool:
    """
    Detect whether the text is technical (CS, math, science) or social science / humanities.
    Returns True if technical.
    """
    technical_indicators = [
        # Programming / CS
        r'\bdef\s+\w+\s*\(', r'\bclass\s+\w+', r'import\s+\w+', r'#.*code',
        r'\bfunction\b', r'\balgorithm\b', r'\barray\b', r'\bloop\b',
        r'\bvariable\b', r'\bsyntax\b', r'\bcompiler\b', r'\bruntime\b',
        r'\bpipeline\b', r'\btoken', r'\bparse', r'\bnode\b', r'\btree\b',
        r'\bstack\b', r'\bqueue\b', r'\brecursi', r'\biterat',
        r'\bO\(n\)', r'\bcomplexity\b', r'\bdataset\b', r'\bmatrix\b',
        # Math
        r'\bequation\b', r'\btheorem\b', r'\bproof\b', r'\bderivative\b',
        r'\bintegral\b', r'\bvector\b', r'\bpolynomial\b', r'\bfunction\b',
        r'[=+\-*/^]{2,}', r'\d+\s*[\+\-\*\/]\s*\d+',
        # Code blocks / output markers
        r'>>>', r'output:', r'input:', r'step\s*\d+:', r'example:',
    ]
    text_lower = text.lower()
    matches = sum(1 for pattern in technical_indicators if re.search(pattern, text_lower))
    return matches >= 3


def _clean_technical_text(text: str) -> str:
    """
    Pre-process technical text before summarisation.
    Removes code output blocks, strips example snippets, and keeps explanatory prose.
    """
    lines = text.split('\n')
    cleaned = []
    skip_next = False

    for line in lines:
        stripped = line.strip()

        # Skip obvious code output lines
        if re.match(r'^(>>>|output|result|example output|in \[\d+\]|out\[\d+\])', stripped, re.IGNORECASE):
            skip_next = True
            continue

        # Skip lines that are purely code (no letters or only symbols)
        if stripped and not re.search(r'[a-zA-Z]{3,}', stripped):
            continue

        # Skip very short lines that are likely labels/headers not useful in summary
        if len(stripped) < 20 and stripped.endswith(':'):
            continue

        if skip_next and stripped == '':
            skip_next = False
            continue

        cleaned.append(line)

    return '\n'.join(cleaned)


def _chunk_text(text: str, chunk_size: int = 3000) -> list:
    """Split long text into overlapping chunks so nothing important is lost."""
    if len(text) <= chunk_size:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        # Try to break at a sentence boundary
        if end < len(text):
            boundary = text.rfind('. ', start, end)
            if boundary > start + chunk_size // 2:
                end = boundary + 1
        chunks.append(text[start:end].strip())
        start = end - 200  # 200 char overlap to avoid cutting mid-concept
    return chunks


def _call_llm(prompt: str, max_tokens: int = 1500) -> Optional[str]:
    """Make a single call to Llama via Hugging Face Inference API."""
    token = getattr(settings, 'HUGGINGFACE_API_TOKEN', '').strip()
    if not token:
        return None
    try:
        response = requests.post(
            "https://router.huggingface.co/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {token}",
                "Content-Type": "application/json",
            },
            json={
                "model": "meta-llama/Llama-3.1-8B-Instruct",
                "messages": [{"role": "user", "content": prompt}],
                "max_tokens": max_tokens,
                "temperature": 0.2,
                "stream": False,
            },
            timeout=120,
        )
        if response.status_code == 200:
            return response.json()["choices"][0]["message"]["content"].strip()
        else:
            print(f"LLM error {response.status_code}: {response.text[:200]}")
    except Exception as e:
        print(f"LLM call error: {e}")
    return None


def _summarize_with_llm(text: str, max_words: int = 800, is_slides: bool = False) -> Optional[str]:
    """
    Generate a genuine understanding-based summary using Llama.
    For long documents, chunks the text and synthesises a final summary.
    """
    if not getattr(settings, 'HUGGINGFACE_API_TOKEN', '').strip():
        return None

    is_technical = _is_technical_content(text)
    chunks = _chunk_text(text, chunk_size=3000)

    if is_slides:
        content_type = "university lecture slide deck where each slide is labelled [Slide N: Title]"
        focus = (
            "Focus on the CONCEPTS and DEFINITIONS being taught in the lecture. "
            "For each major topic or step, explain what it is, how it works, and why it matters. "
            "Do NOT reproduce code — explain in plain English what it demonstrates. "
            "Ignore analogies and examples used to illustrate points; explain the underlying concept instead."
        )
    elif is_technical:
        content_type = "technical study material (computer science, mathematics, or engineering)"
        focus = (
            "Focus on explaining the key concepts, algorithms, processes, and definitions covered. "
            "Explain technical terms in plain academic English. "
            "Do NOT reproduce code snippets — describe what they do and why. "
            "Explain each step or procedure in terms of its purpose and outcome."
        )
    else:
        content_type = "academic study material"
        focus = (
            "Identify and explain the main arguments, theories, concepts, and conclusions. "
            "Show how the ideas relate to each other. "
            "Write as if explaining to a student who needs to understand the topic deeply, not just recall facts."
        )

    # For short documents: single call
    if len(chunks) == 1:
        prompt = f"""You are an expert academic tutor. Read the {content_type} below and write a thorough, understanding-based summary.

{focus}

Write a well-structured summary of approximately {max_words} words. Use clear paragraphs. Do NOT use bullet points. Write in full academic prose as if explaining this topic to a fellow student.

MATERIAL:
{chunks[0]}

Write the summary now:"""
        return _call_llm(prompt, max_tokens=2000)

    # For long documents: summarise each chunk then synthesise
    chunk_summaries = []
    for i, chunk in enumerate(chunks):
        chunk_prompt = f"""You are an expert academic tutor. Read this excerpt (part {i+1} of {len(chunks)}) from a {content_type} and write a concise but thorough explanation of what it covers.

{focus}

Write 2-4 paragraphs in plain academic prose. Do NOT use bullet points.

EXCERPT:
{chunk}

Write the explanation now:"""
        result = _call_llm(chunk_prompt, max_tokens=800)
        if result:
            chunk_summaries.append(result)

    if not chunk_summaries:
        return None

    # Synthesise chunk summaries into a final coherent summary
    combined = "\n\n".join(chunk_summaries)
    synthesis_prompt = f"""You are an expert academic tutor. Below are notes covering different parts of a {content_type}. Combine them into one coherent, well-structured summary of approximately {max_words} words.

Rules:
- Write in flowing academic prose with clear paragraphs
- Do NOT use bullet points
- Remove any repetition between sections
- Make sure the summary reads as one unified explanation, not separate parts
- Explain concepts so a student can genuinely understand the material

NOTES:
{combined[:4000]}

Write the final unified summary now:"""

    return _call_llm(synthesis_prompt, max_tokens=2000)


def summarize_text(text: str, min_words: int = 300, max_words: int = 800) -> str:
    """
    Generate a summary using the LLM API.
    Raises an error message string if no API key is configured or the call fails.
    """
    if not text or not text.strip():
        raise ValueError("No content could be extracted from the uploaded file.")

    is_technical = _is_technical_content(text)
    is_slides = bool(re.search(r'\[Slide \d+:', text))

    hf_key = getattr(settings, 'HUGGINGFACE_API_TOKEN', '').strip()

    if not hf_key:
        raise ValueError(
            "No AI API key is configured. Please add GEMINI_API_KEY or "
            "HUGGINGFACE_API_TOKEN to your .env file and restart the server."
        )

    summary = _summarize_with_llm(text, max_words=max_words, is_slides=is_slides)

    if not summary or not summary.strip():
        raise ValueError(
            "The AI could not generate a summary for this material. "
            "This may be due to a network issue or API error. Please try again."
        )

    return summary




def expand_summary(summary: str, original_text: str, min_words: int, max_words: int) -> str:
    """
    Expand a summary to meet minimum word count by adding more content.
    
    Args:
        summary: Current summary
        original_text: Original full text
        min_words: Minimum word count
        max_words: Maximum word count
    
    Returns:
        Expanded summary
    """
    summary_words = count_words(summary)
    if summary_words >= min_words:
        return summary
    
    # Get sentences from original text not in summary
    original_sentences = re.split(r'(?<=[.!?])\s+', original_text)
    original_sentences = [s.strip() for s in original_sentences if len(s.strip()) > 15]
    
    summary_sentences = re.split(r'(?<=[.!?])\s+', summary)
    summary_sentences = [s.strip() for s in summary_sentences]
    
    # Add sentences until we reach min_words
    current_words = summary_words
    additional_sentences = []
    
    for sentence in original_sentences:
        if sentence not in summary_sentences:
            sentence_words = count_words(sentence)
            if current_words + sentence_words <= max_words:
                additional_sentences.append(sentence)
                current_words += sentence_words
            if current_words >= min_words:
                break
    
    # Combine summary with additional sentences
    expanded = summary + ' ' + ' '.join(additional_sentences)
    
    # If still too short, just take from the beginning of original text
    if count_words(expanded) < min_words:
        words_needed = min_words - count_words(expanded)
        extra_text = ' '.join(original_text.split()[:words_needed])
        expanded = expanded + ' ' + extra_text
    
    return expanded.strip()


def extract_key_concepts(text: str, num_concepts: int = 15) -> List[str]:
    """
    Extract key concepts/keywords from text.
    
    Args:
        text: The text to analyze
        num_concepts: Number of key concepts to extract
    
    Returns:
        List of key concepts
    """
    # Clean and tokenize - include multi-word phrases
    words = re.findall(r'\b[A-Za-z][a-zA-Z]*(?:\s+[A-Za-z][a-zA-Z]*){0,2}\b', text)
    
    # Common stop words
    stop_words = {
        'the', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been', 'being',
        'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could',
        'should', 'may', 'might', 'must', 'shall', 'can', 'need', 'dare',
        'ought', 'used', 'to', 'of', 'in', 'for', 'on', 'with', 'at', 'by',
        'from', 'as', 'into', 'through', 'during', 'before', 'after', 'above',
        'below', 'between', 'under', 'and', 'but', 'or', 'yet', 'so', 'if',
        'because', 'although', 'though', 'while', 'where', 'when', 'that',
        'which', 'who', 'whom', 'whose', 'what', 'this', 'these', 'those',
        'i', 'you', 'he', 'she', 'it', 'we', 'they', 'me', 'him', 'her',
        'us', 'them', 'my', 'your', 'his', 'its', 'our', 'their', 'mine',
        'yours', 'hers', 'ours', 'theirs', 'myself', 'yourself', 'himself',
        'herself', 'itself', 'ourselves', 'yourselves', 'themselves',
        'very', 'really', 'just', 'now', 'then', 'than', 'only', 'also',
        'well', 'even', 'more', 'most', 'some', 'any', 'all', 'each',
        'every', 'both', 'either', 'neither', 'one', 'two', 'first',
        'last', 'next', 'other', 'another', 'such', 'same', 'different'
    }
    
    # Count word frequencies (excluding stop words)
    word_freq = {}
    for word in words:
        word_lower = word.lower()
        if word_lower not in stop_words and len(word) > 2:
            # Boost score for capitalized words (likely proper nouns/important terms)
            boost = 1.5 if word[0].isupper() else 1.0
            word_freq[word] = word_freq.get(word, 0) + boost
    
    # Get top concepts
    sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
    key_concepts = [word for word, freq in sorted_words[:num_concepts]]
    
    return key_concepts


def generate_quiz_questions(
    text: str, 
    key_concepts: List[str], 
    num_questions: int = 5,
    question_types: List[str] = None
) -> List[Dict]:
    """
    Generate smart quiz questions from text using the SmartQuizGenerator.
    
    Args:
        text: The source text
        key_concepts: List of key concepts (kept for compatibility)
        num_questions: Number of questions to generate
        question_types: Types of questions to generate
    
    Returns:
        List of question dictionaries
    """
    # Import here to avoid circular imports
    from .quiz_generator import generate_smart_quiz_questions
    
    return generate_smart_quiz_questions(text, num_questions, question_types)


def process_study_material(text: str) -> Tuple[str, List[str]]:
    """
    Process study material to generate summary and key concepts.
    
    Args:
        text: The study material text
    
    Returns:
        Tuple of (summary, key_concepts)
    """
    summary = summarize_text(text, min_words=300, max_words=800)
    key_concepts = extract_key_concepts(text, num_concepts=15)
    return summary, key_concepts
